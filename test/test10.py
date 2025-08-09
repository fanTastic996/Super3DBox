from pptx import Presentation
from pptx.util import Inches
import os

def images_to_ppt(image_folder, output_file="output.pptx"):
    prs = Presentation()  # 创建PPT对象
    # 遍历图片文件夹，过滤格式
    image_exts = ['.jpg', '.jpeg', '.png', '.gif']
    image_files = [f for f in os.listdir(image_folder) if os.path.splitext(f)[1].lower() in image_exts]
    image_files = sorted(image_files, key=lambda x: x.split('.')[0])
    for img_file in image_files:
        img_path = os.path.join(image_folder, img_file)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 使用空白布局（代码6）
        left = Inches(1)   # 图片左边界（1英寸）
        top = Inches(1)  # 图片上边界（1.5英寸）
        slide.shapes.add_picture(img_path, left, top, height=Inches(6))  # 固定高度，宽度按比例缩放
    prs.save(output_file)
    print(f"生成成功！共插入 {len(image_files)} 张图片，保存至：{output_file}")

# 使用示例 ↓
images_to_ppt(image_folder="/home/lyq/Downloads/chinasi/", output_file="chinasi.pptx")